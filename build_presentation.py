"""
Gera a apresentação PowerPoint completa do projeto Praso Credit Risk.
"""

import io
import json
import warnings
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Caminhos ─────────────────────────────────────────────────────────────────
ROOT      = Path(__file__).resolve().parent
DATA_DIR  = ROOT / 'data'
OUT_PATH  = ROOT / 'Praso_CreditRisk_Apresentacao.pptx'

# ── Paleta ───────────────────────────────────────────────────────────────────
def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

DARK   = rgb('#263238')
WHITE  = rgb('#FFFFFF')
LGRAY  = rgb('#ECEFF1')
MGRAY  = rgb('#546E7A')
ORANGE = rgb('#E65100')
BLUE   = rgb('#1565C0')
PURPLE = rgb('#4527A0')
GREEN  = rgb('#2E7D32')
RED    = rgb('#C62828')
TEAL   = rgb('#00695C')
AMBER  = rgb('#FF8F00')
LBLUE  = rgb('#1E88E5')   # tier behavioral
LORANG = rgb('#FB8C00')   # tier application
LPURP  = rgb('#8E24AA')   # tier manual

# ── Helpers ──────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height
BLANK = prs.slide_layouts[6]

def add_rect(slide, l, t, w, h, color, line=False, line_color=None, lw=None):
    shp = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line and line_color:
        shp.line.color.rgb = line_color
        if lw:
            shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    return shp

def add_text(slide, text, l, t, w, h,
             size=13, bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, wrap=True):
    color = color or DARK
    txb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_bullets(slide, items, l, t, w, h, size=13, color=None,
                bullet='▸ ', spacing=6, bold_first=False):
    color = color or DARK
    txb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
    return txb

def add_slide_header(title, accent, subtitle=''):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, LGRAY)
    add_rect(slide, 0, 0, W, Inches(1.2), accent)
    add_rect(slide, 0, Inches(1.2), W, Pt(4), accent)
    add_text(slide, title,
             Inches(0.45), Inches(0.1), Inches(12.4), Inches(0.8),
             size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.45), Inches(0.82), Inches(12.4), Inches(0.35),
                 size=11, color=rgb('#B0BEC5'))
    return slide

def add_slide_number(slide, n):
    add_text(slide, str(n),
             Inches(12.6), Inches(7.1), Inches(0.6), Inches(0.3),
             size=9, color=MGRAY, align=PP_ALIGN.RIGHT)

def fig_to_buf(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf

def add_image(slide, buf_or_path, l, t, w, h):
    slide.shapes.add_picture(buf_or_path if hasattr(buf_or_path, 'read')
                             else str(buf_or_path), int(l), int(t), int(w), int(h))

def kpi_box(slide, l, t, w, h, value, label, color):
    add_rect(slide, l, t, w, h, color)
    add_text(slide, value,
             l, t + Inches(0.1), w, Inches(0.65),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label,
             l, t + Inches(0.68), w, Inches(0.42),
             size=10, color=rgb('#CFD8DC'), align=PP_ALIGN.CENTER)

def section_box(slide, l, t, w, h, title, body, color):
    add_rect(slide, l, t, w, h, color)
    mid = (h - Inches(0.55)) / 2
    add_text(slide, title,
             l + Inches(0.15), t + Inches(0.1), w - Inches(0.3), Inches(0.4),
             size=11, bold=True, color=WHITE)
    add_text(slide, body,
             l + Inches(0.15), t + Inches(0.45), w - Inches(0.3), h - Inches(0.5),
             size=9.5, color=rgb('#E3F2FD'))

# ─────────────────────────────────────────────────────────────────────────────
# CHARTS
# ─────────────────────────────────────────────────────────────────────────────

def chart_data_overview():
    fig, axes = plt.subplots(1, 2, figsize=(10, 4), facecolor='#FAFAFA')
    # Left: pie - clients with/without behavioral data
    sizes  = [664, 2336]
    labels = ['Com histórico\ncomportamental\n(664)', 'Sem histórico\n(2.336)']
    colors = ['#1565C0', '#90A4AE']
    axes[0].pie(sizes, labels=labels, colors=colors, autopct='%1.1f%%',
                startangle=140, textprops={'fontsize': 10})
    axes[0].set_title('Distribuição dos Clientes', fontsize=12, fontweight='bold', pad=8)

    # Right: default rate overall + per group
    categories = ['Portfolio\ntotal', 'Com dados\ncomportamentais', 'Sem dados\ncomportamentais']
    rates      = [31.3, 20.8, 33.9]
    bar_colors = ['#546E7A', '#1565C0', '#FB8C00']
    bars = axes[1].bar(categories, rates, color=bar_colors, width=0.5, edgecolor='white')
    for bar, r in zip(bars, rates):
        axes[1].text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.5,
                     f'{r:.1f}%', ha='center', va='bottom', fontsize=10, fontweight='bold')
    axes[1].set_ylim(0, 42)
    axes[1].set_ylabel('Taxa de Inadimplência (%)', fontsize=10)
    axes[1].set_title('Taxa de Inadimplência por Grupo', fontsize=12, fontweight='bold')
    axes[1].axhline(31.3, color='#C62828', ls='--', lw=1.2, label='Média geral (31.3%)')
    axes[1].legend(fontsize=9)
    fig.tight_layout(pad=2)
    return fig_to_buf(fig)

def chart_model_auc():
    fig, ax = plt.subplots(figsize=(10, 4.5), facecolor='#FAFAFA')
    models = ['App Model\n(Baseline)', 'App Model\n(Tuned XGBoost)',
              'Behavioral\n(RF Tuned)', 'Sistema\n(Rule-based)', 'Sistema\n(ML Router)']
    aucs   = [0.726, 0.772, 0.770, 0.850, 0.848]
    colors = ['#90A4AE', '#E65100', '#1565C0', '#4527A0', '#6A1B9A']
    bars   = ax.bar(models, aucs, color=colors, width=0.55, edgecolor='white')
    for bar, v in zip(bars, aucs):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.002,
                f'{v:.3f}', ha='center', va='bottom', fontsize=11, fontweight='bold')
    ax.set_ylim(0.65, 0.90)
    ax.set_ylabel('ROC-AUC', fontsize=11)
    ax.set_title('Comparação de Modelos — ROC-AUC', fontsize=13, fontweight='bold')
    ax.axhline(0.726, color='#90A4AE', ls='--', lw=1, label='Baseline (0.726)')
    ax.legend(fontsize=9)
    ax.spines[['top','right']].set_visible(False)
    ax.set_facecolor('#F5F5F5')
    fig.tight_layout()
    return fig_to_buf(fig)

def chart_routing_dist():
    fig, axes = plt.subplots(1, 2, figsize=(10, 4), facecolor='#FAFAFA')

    # Pie
    tiers  = ['BEHAVIORAL\n(664)', 'APPLICATION\n(998)', 'MANUAL_REVIEW\n(1.338)']
    counts = [664, 998, 1338]
    colors = ['#1E88E5', '#FB8C00', '#8E24AA']
    axes[0].pie(counts, labels=tiers, colors=colors, autopct='%1.1f%%',
                startangle=140, textprops={'fontsize': 9.5})
    axes[0].set_title('Distribuição de Roteamento (3.000)', fontsize=11, fontweight='bold')

    # Default rate per tier
    tier_names = ['BEHAVIORAL', 'APPLICATION', 'MANUAL_REVIEW']
    dr = [20.8, 33.4, 35.1]
    axes[1].barh(tier_names, dr, color=colors, edgecolor='white', height=0.5)
    for i, v in enumerate(dr):
        axes[1].text(v + 0.3, i, f'{v:.1f}%', va='center', fontsize=11, fontweight='bold')
    axes[1].set_xlim(0, 45)
    axes[1].set_xlabel('Taxa de Inadimplência (%)', fontsize=10)
    axes[1].set_title('Taxa de Inadimplência por Tier', fontsize=11, fontweight='bold')
    axes[1].axvline(31.3, color='#C62828', ls='--', lw=1.2, label='Média (31.3%)')
    axes[1].legend(fontsize=9)
    axes[1].spines[['top','right']].set_visible(False)
    axes[1].set_facecolor('#F5F5F5')
    fig.tight_layout(pad=2)
    return fig_to_buf(fig)

def chart_shap():
    fig, axes = plt.subplots(1, 2, figsize=(11, 4.5), facecolor='#FAFAFA')

    # App Model
    app_feats = ['idade_cnpj_mid', 'serasa_socio_\ntem_negativacao',
                 'serasa_n_setores', 'capital_social_mid', 'google_maps_\navaliacao_mid']
    app_vals  = [0.553, 0.370, 0.193, 0.151, 0.118]
    y = np.arange(len(app_feats))
    axes[0].barh(y, app_vals, color='#E65100', edgecolor='white', height=0.6)
    axes[0].set_yticks(y)
    axes[0].set_yticklabels(app_feats, fontsize=9)
    axes[0].set_xlabel('|SHAP| médio', fontsize=10)
    axes[0].set_title('Modelo de Aplicação\n(XGBoost)', fontsize=11, fontweight='bold', color='#E65100')
    axes[0].spines[['top','right']].set_visible(False)
    axes[0].set_facecolor('#FFF3E0')
    for i, v in enumerate(app_vals):
        axes[0].text(v + 0.005, i, f'{v:.3f}', va='center', fontsize=9)

    # Behavioral Model
    beh_feats = ['idade_cnpj_mid', 'delay_mean', 'orders_per_month',
                 'pct_orders_delayed', 'log_valor_mean']
    beh_vals  = [0.038, 0.037, 0.037, 0.031, 0.029]
    axes[1].barh(y, beh_vals, color='#1565C0', edgecolor='white', height=0.6)
    axes[1].set_yticks(y)
    axes[1].set_yticklabels(beh_feats, fontsize=9)
    axes[1].set_xlabel('|SHAP| médio', fontsize=10)
    axes[1].set_title('Modelo Comportamental\n(Random Forest)', fontsize=11, fontweight='bold', color='#1565C0')
    axes[1].spines[['top','right']].set_visible(False)
    axes[1].set_facecolor('#E3F2FD')
    for i, v in enumerate(beh_vals):
        axes[1].text(v + 0.0005, i, f'{v:.3f}', va='center', fontsize=9)

    fig.suptitle('Top-5 Features — Importância SHAP', fontsize=13, fontweight='bold', y=1.01)
    fig.tight_layout(pad=2)
    return fig_to_buf(fig)

def chart_scenarios():
    fig, axes = plt.subplots(1, 2, figsize=(11, 4.5), facecolor='#FAFAFA')
    scenarios = ['Baseline\n(App, p<0.25)', 'Router\nRecomendado', 'Router\nConservador', 'Router\nPermissivo']
    approval  = [21.4, 39.0, 21.4, 60.7]
    default   = [5.6,   7.3,  2.8, 11.7]
    loss      = [0.43,  1.03, 0.22,  2.56]
    x    = np.arange(len(scenarios))
    cols = ['#90A4AE', '#4527A0', '#1565C0', '#E65100']

    bars1 = axes[0].bar(x, approval, color=cols, width=0.55, edgecolor='white')
    for bar, v in zip(bars1, approval):
        axes[0].text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.5,
                     f'{v:.1f}%', ha='center', va='bottom', fontsize=10, fontweight='bold')
    axes[0].set_xticks(x)
    axes[0].set_xticklabels(scenarios, fontsize=9)
    axes[0].set_ylabel('Taxa de Aprovação (%)', fontsize=10)
    axes[0].set_title('Taxa de Aprovação por Cenário', fontsize=11, fontweight='bold')
    axes[0].set_ylim(0, 75)
    axes[0].spines[['top','right']].set_visible(False)
    axes[0].set_facecolor('#F5F5F5')

    bars2 = axes[1].bar(x, loss, color=cols, width=0.55, edgecolor='white')
    for bar, v in zip(bars2, loss):
        axes[1].text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.02,
                     f'R${v:.2f}M', ha='center', va='bottom', fontsize=10, fontweight='bold')
    axes[1].set_xticks(x)
    axes[1].set_xticklabels(scenarios, fontsize=9)
    axes[1].set_ylabel('Perda Esperada (R$M)', fontsize=10)
    axes[1].set_title('Perda Esperada por Cenário', fontsize=11, fontweight='bold')
    axes[1].set_ylim(0, 3.2)
    axes[1].spines[['top','right']].set_visible(False)
    axes[1].set_facecolor('#F5F5F5')

    fig.tight_layout(pad=2)
    return fig_to_buf(fig)

# ─────────────────────────────────────────────────────────────────────────────
# PRE-GENERATE CHARTS
# ─────────────────────────────────────────────────────────────────────────────
print('Gerando gráficos...')
img_data_overview = chart_data_overview()
img_model_auc     = chart_model_auc()
img_routing       = chart_routing_dist()
img_shap          = chart_shap()
img_scenarios     = chart_scenarios()
arch_path         = DATA_DIR / '00_system_architecture.png'
print('Gráficos prontos.')

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Capa
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, DARK)
add_rect(slide, 0, Inches(5.55), W, Inches(1.95), PURPLE)
# Accent line
add_rect(slide, 0, Inches(5.55), Pt(8), Inches(1.95), ORANGE)

add_text(slide, 'Praso Credit Risk',
         Inches(0.7), Inches(1.2), Inches(11.5), Inches(1.2),
         size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 'Previsão de Inadimplência com Model Routing',
         Inches(0.7), Inches(2.4), Inches(11.5), Inches(0.7),
         size=20, color=rgb('#B0BEC5'), align=PP_ALIGN.CENTER)
add_text(slide, 'Arquitetura inspirada em Nubank · XGBoost + Random Forest + Router (Score A / Score B)',
         Inches(0.7), Inches(3.1), Inches(11.5), Inches(0.5),
         size=12, color=rgb('#78909C'), align=PP_ALIGN.CENTER, italic=True)

# Metrics strip
kpi_data = [
    ('3.000', 'Clientes PME', BLUE),
    ('31,3%', 'Taxa inadimplência', RED),
    ('0,850', 'AUC do sistema', PURPLE),
    ('+4,6pp', 'Lift vs. baseline', GREEN),
    ('39,0%', 'Taxa de aprovação', TEAL),
]
for i, (v, l, c) in enumerate(kpi_data):
    x = Inches(0.5 + i * 2.48)
    add_rect(slide, x, Inches(5.65), Inches(2.3), Inches(1.7), c)
    add_text(slide, v, x, Inches(5.75), Inches(2.3), Inches(0.75),
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, l, x, Inches(6.42), Inches(2.3), Inches(0.5),
             size=9, color=rgb('#E3F2FD'), align=PP_ALIGN.CENTER)

add_text(slide, 'Projeto Universitário de Machine Learning · 2026',
         Inches(0.5), Inches(7.15), Inches(12.0), Inches(0.3),
         size=9, color=rgb('#546E7A'), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Agenda
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide_header('Agenda', DARK)
add_slide_number(slide, 2)

agenda = [
    ('01', 'O Problema',              ORANGE),
    ('02', 'Os Dados',                TEAL),
    ('03', 'Engenharia de Features',  TEAL),
    ('04', 'Modelo de Aplicação',     ORANGE),
    ('05', 'Modelo Comportamental',   BLUE),
    ('06', 'Model Routing',           PURPLE),
    ('07', 'Interpretabilidade',      MGRAY),
    ('08', 'Política de Crédito',     GREEN),
    ('09', 'Resultados & Próximos Passos', RED),
]
cols_per_row = 3
for i, (num, topic, color) in enumerate(agenda):
    row = i // cols_per_row
    col = i  % cols_per_row
    x = Inches(0.5 + col * 4.27)
    y = Inches(1.4 + row * 1.8)
    add_rect(slide, x, y, Inches(3.9), Inches(1.5), color)
    add_text(slide, num, x + Inches(0.15), y + Inches(0.1), Inches(0.7), Inches(0.5),
             size=22, bold=True, color=rgb('#FFFFFF80'))
    add_text(slide, topic, x + Inches(0.15), y + Inches(0.55), Inches(3.6), Inches(0.75),
             size=13, bold=True, color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — O Problema
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide_header('O Problema', ORANGE, 'Risco de crédito para PMEs no marketplace Praso')
add_slide_number(slide, 3)

add_rect(slide, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.8), ORANGE)
add_text(slide, 'Contexto',
         Inches(0.55), Inches(1.5), Inches(5.5), Inches(0.45),
         size=14, bold=True, color=WHITE)
add_bullets(slide,
    ['Praso é um marketplace B2B que oferece crédito a pequenas e médias empresas (PMEs)',
     'Risco de crédito afeta diretamente a saúde financeira da plataforma',
     '31,3% dos clientes são inadimplentes — acima da média do setor',
     'Decisões manuais são lentas, custosas e inconsistentes',
     'A heterogeneidade dos clientes (com/sem histórico) dificulta um único modelo'],
    Inches(0.55), Inches(1.95), Inches(5.5), Inches(5.0),
    size=11.5, color=WHITE, bullet='• ', spacing=10)

add_rect(slide, Inches(6.5), Inches(1.4), Inches(6.5), Inches(2.6), BLUE)
add_text(slide, 'Objetivo',
         Inches(6.65), Inches(1.5), Inches(6.2), Inches(0.45),
         size=14, bold=True, color=WHITE)
add_bullets(slide,
    ['Prever inadimplência antes de conceder crédito',
     'Usar dados de aplicação E histórico comportamental (pedidos)',
     'Rotear cada cliente ao modelo mais adequado',
     'Gerar uma política de crédito operacional e auditável'],
    Inches(6.65), Inches(1.95), Inches(6.2), Inches(1.9),
    size=11.5, color=WHITE, bullet='▸ ', spacing=8)

add_rect(slide, Inches(6.5), Inches(4.2), Inches(6.5), Inches(3.0), TEAL)
add_text(slide, 'Abordagem — Model Routing (inspirado Nubank)',
         Inches(6.65), Inches(4.3), Inches(6.2), Inches(0.45),
         size=13, bold=True, color=WHITE)
add_bullets(slide,
    ['Score A (riqueza de dados) avalia o volume de informação disponível',
     'Score B (complexidade de perfil) mede incerteza do modelo de aplicação',
     'Router seleciona o modelo ótimo para cada cliente',
     'Resultado: sistema AUC 0,850 vs baseline 0,726 (+17%)'],
    Inches(6.65), Inches(4.8), Inches(6.2), Inches(2.3),
    size=11, color=WHITE, bullet='▸ ', spacing=8)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Os Dados
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide_header('Os Dados', TEAL, 'Duas fontes de dados · 3.000 clientes PME')
add_slide_number(slide, 4)

# Two source boxes
for i, (title, body, color) in enumerate([
    ('aplicacao.csv',
     '3.000 linhas × 28 colunas\n\nDados cadastrais: capital social,\nidade do CNPJ, CNAE, avaliação\nGoogle Maps, dados Serasa (sócios,\nnegativos, credores), segmento',
     ORANGE),
    ('pedidos.csv',
     '664 clientes com histórico\n\nDados de pedidos: valor médio,\nmáximo, mínimo e desvio,\ntaxa de atraso, dias de atraso\nmédio e máximo, pedidos/mês',
     BLUE),
]):
    x = Inches(0.4 + i * 4.6)
    add_rect(slide, x, Inches(1.4), Inches(4.2), Inches(2.7), color)
    add_text(slide, title, x + Inches(0.15), Inches(1.5), Inches(3.9), Inches(0.45),
             size=13, bold=True, color=WHITE)
    add_text(slide, body, x + Inches(0.15), Inches(1.95), Inches(3.9), Inches(2.1),
             size=10.5, color=rgb('#E3F2FD'))

add_image(slide, img_data_overview, Inches(9.1), Inches(1.25), Inches(4.1), Inches(4.0))

# Stats row
stats = [
    ('3.000', 'clientes total', MGRAY),
    ('664',   'com dados\ncomportamentais', BLUE),
    ('31,3%', 'taxa de\ninadimplência', RED),
    ('22,1%', 'clientes com\nhistórico', TEAL),
]
for i, (v, l, c) in enumerate(stats):
    x = Inches(0.4 + i * 2.2)
    add_rect(slide, x, Inches(4.35), Inches(2.0), Inches(1.4), c)
    add_text(slide, v, x, Inches(4.42), Inches(2.0), Inches(0.65),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, l, x, Inches(5.0), Inches(2.0), Inches(0.7),
             size=9, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.4), Inches(5.95), Inches(8.7), Inches(1.3), rgb('#263238'))
add_text(slide, '⚠  Desafio principal: dois segmentos de clientes com informação muito diferente.',
         Inches(0.55), Inches(5.98), Inches(8.4), Inches(0.5),
         size=12, bold=True, color=rgb('#FFF176'))
add_text(slide, 'Clientes SEM histórico de pedidos precisam de um modelo mais conservador; '
         'clientes COM histórico permitem um modelo mais granular e preciso.',
         Inches(0.55), Inches(6.42), Inches(8.4), Inches(0.75),
         size=10.5, color=rgb('#B0BEC5'))

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Engenharia de Features
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide_header('Engenharia de Features', TEAL,
                         'Transformação de dados brutos em representações numéricas para os modelos')
add_slide_number(slide, 5)

feat_groups = [
    ('Intervalos → Midpoint',
     'capital_social, idade_cnpj,\ngoogle_maps_avaliacao\n→ extrai ponto médio numérico',
     TEAL),
    ('CNAE → Divisão',
     'cnae_codigo → cnae_divisao\n→ agrupa por setor econômico\n(2 primeiros dígitos)',
     TEAL),
    ('Serasa — Credores',
     'serasa_credores (texto)\n→ serasa_n_setores\n→ contagem de setores devedores',
     ORANGE),
    ('Features Comportamentais',
     'log_valor_* (log1p das estatísticas\nde valor), has_any_delay,\ndelay_spike_ratio',
     BLUE),
    ('Encoding Booleano',
     'Colunas bool → int\npara compatibilidade com\nsklearn pipelines',
     MGRAY),
    ('Pipeline ColumnTransformer',
     'remainder="drop": aceita\nsuperconjunto de colunas\nde forma segura',
     PURPLE),
]

for i, (title, body, color) in enumerate(feat_groups):
    row = i // 3
    col = i  % 3
    x = Inches(0.4 + col * 4.27)
    y = Inches(1.45 + row * 2.8)
    add_rect(slide, x, y, Inches(4.0), Inches(2.5), color)
    add_text(slide, title, x + Inches(0.15), y + Inches(0.1), Inches(3.7), Inches(0.45),
             size=11.5, bold=True, color=WHITE)
    add_text(slide, body, x + Inches(0.15), y + Inches(0.55), Inches(3.7), Inches(1.8),
             size=10.5, color=rgb('#E3F2FD'))

add_rect(slide, Inches(0.4), Inches(7.08), Inches(12.5), Inches(0.28), DARK)
add_text(slide, 'Resultado: 41 features para Modelo de Aplicação  ·  53 features para Modelo Comportamental',
         Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.28),
         size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Modelo de Aplicação
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide_header('Modelo de Aplicação', ORANGE, 'XGBoost · Treinado em 3.000 clientes')
add_slide_number(slide, 6)

add_rect(slide, Inches(0.4), Inches(1.4), Inches(4.5), Inches(5.8), ORANGE)
add_text(slide, 'Resultado',
         Inches(0.55), Inches(1.5), Inches(4.2), Inches(0.4),
         size=13, bold=True, color=WHITE)
add_bullets(slide, [
    'Algoritmo: XGBoost (Gradient Boosting)',
    'ROC-AUC: 0,772 (tuned vs 0,726 baseline)',
    'Average Precision: 0,637',
    'Recall @ threshold 0,30: 0,72',
    'Treinamento: 70% (2.100) / Teste: 30% (900)',
    '',
    'Tuning via RandomizedSearchCV',
    '100 iterações · 5-fold CV',
    'Parâmetros: n_estimators, max_depth,\nlearning_rate, subsample, colsample',
], Inches(0.55), Inches(1.9), Inches(4.15), Inches(5.0),
   size=10.5, color=WHITE, bullet='▸ ', spacing=7)

add_image(slide, img_model_auc, Inches(5.1), Inches(1.35), Inches(7.9), Inches(3.8))

add_rect(slide, Inches(5.1), Inches(5.35), Inches(7.9), Inches(1.9), DARK)
add_text(slide, 'Principais Inputs (por SHAP)',
         Inches(5.25), Inches(5.42), Inches(7.6), Inches(0.4),
         size=12, bold=True, color=WHITE)
for i, (feat, val) in enumerate([
    ('idade_cnpj_mid', '0,553'),
    ('serasa_socio_tem_negativacao', '0,370'),
    ('serasa_n_setores', '0,193'),
]):
    x = Inches(5.25 + i * 2.6)
    add_rect(slide, x, Inches(5.85), Inches(2.4), Inches(1.2), ORANGE)
    add_text(slide, feat, x + Inches(0.1), Inches(5.88), Inches(2.2), Inches(0.55),
             size=9.5, bold=True, color=WHITE)
    add_text(slide, f'|SHAP| = {val}', x + Inches(0.1), Inches(6.35), Inches(2.2), Inches(0.45),
             size=11, color=rgb('#FFE0B2'), bold=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Modelo Comportamental
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide_header('Modelo Comportamental', BLUE, 'Random Forest · Treinado em 664 clientes com histórico de pedidos')
add_slide_number(slide, 7)

add_rect(slide, Inches(0.4), Inches(1.4), Inches(4.5), Inches(5.8), BLUE)
add_text(slide, 'Resultado',
         Inches(0.55), Inches(1.5), Inches(4.2), Inches(0.4),
         size=13, bold=True, color=WHITE)
add_bullets(slide, [
    'Algoritmo: Random Forest',
    'ROC-AUC: 0,770 (vs 0,623 baseline app-only)',
    'Lift sobre baseline: +0,1467',
    'Average Precision: 0,594',
    '',
    'Subconjunto: 664 clientes com pedidos',
    'Taxa inadimplência: 20,8% (vs 31,3% geral)',
    '',
    'Features exclusivas comportamentais:',
    '   delay_mean, delay_max, log_valor_*,',
    '   orders_per_month, pct_orders_delayed,',
    '   has_any_delay, delay_spike_ratio',
], Inches(0.55), Inches(1.9), Inches(4.15), Inches(5.0),
   size=10.5, color=WHITE, bullet='▸ ', spacing=6)

add_rect(slide, Inches(5.1), Inches(1.4), Inches(7.9), Inches(5.8), rgb('#E3F2FD'))
add_text(slide, 'Por que Random Forest?',
         Inches(5.25), Inches(1.5), Inches(7.6), Inches(0.4),
         size=12, bold=True, color=BLUE)
add_bullets(slide, [
    'Amostra menor (664): RF menos propenso a overfitting que XGBoost',
    'Interpreta bem features contínuas de distribuição irregular (log_valor_*)',
    'Suporte nativo a features mistas (numéricas + binárias)',
    'SHAP com TreeExplainer funciona diretamente em Random Forest',
    'Validação: AUC estável entre folds (σ = 0,02)',
], Inches(5.25), Inches(1.95), Inches(7.6), Inches(2.2),
   size=11, color=DARK, bullet='▸ ')

add_rect(slide, Inches(5.1), Inches(4.3), Inches(3.75), Inches(2.7), BLUE)
add_text(slide, 'Comparação no subconjunto (664)',
         Inches(5.25), Inches(4.4), Inches(3.45), Inches(0.4),
         size=11, bold=True, color=WHITE)
for i, (model, auc) in enumerate([
    ('App Model (baseline)', 0.623),
    ('Behavioral RF (tuned)', 0.770),
]):
    y = Inches(4.9 + i * 0.85)
    add_rect(slide, Inches(5.25), y, Inches(3.45), Inches(0.7),
             GREEN if i == 1 else rgb('#546E7A'))
    add_text(slide, f'{model}\nAUC = {auc:.3f}',
             Inches(5.4), y + Inches(0.05), Inches(3.15), Inches(0.6),
             size=10.5, bold=(i==1), color=WHITE)

add_rect(slide, Inches(9.05), Inches(4.3), Inches(3.95), Inches(2.7), rgb('#1A237E'))
add_text(slide, 'Top Features (SHAP)',
         Inches(9.2), Inches(4.4), Inches(3.65), Inches(0.4),
         size=11, bold=True, color=WHITE)
for i, (f, v) in enumerate([
    ('idade_cnpj_mid', 0.038), ('delay_mean', 0.037),
    ('orders_per_month', 0.037), ('pct_orders_delayed', 0.031),
]):
    y = Inches(4.85 + i * 0.5)
    bar_w = Inches(v / 0.040 * 2.8)
    add_rect(slide, Inches(9.2), y, bar_w, Inches(0.38), LBLUE)
    add_text(slide, f'{f}  {v:.3f}',
             Inches(9.2), y + Inches(0.04), Inches(3.6), Inches(0.35),
             size=9, color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Arquitetura do Sistema (diagrama completo)
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide_header('Arquitetura do Sistema', DARK,
                         'Visão completa: ingestão de dados → feature engineering → modelos → router → política')
add_slide_number(slide, 8)

if arch_path.exists():
    add_image(slide, arch_path, Inches(0.3), Inches(1.25), Inches(12.7), Inches(6.1))
else:
    add_text(slide, 'Diagrama não encontrado. Execute o script de arquitetura primeiro.',
             Inches(1), Inches(3), Inches(11), Inches(1),
             size=14, color=RED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Model Routing
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide_header('Model Routing — Arquitetura', PURPLE,
                         'Roteamento baseado em pontuação (Score A + Score B) · Inspirado em Nubank')
add_slide_number(slide, 9)

# Score A box
add_rect(slide, Inches(0.4), Inches(1.4), Inches(6.0), Inches(2.55), PURPLE)
add_text(slide, 'Score A — Riqueza de Dados',
         Inches(0.55), Inches(1.5), Inches(5.7), Inches(0.45),
         size=13, bold=True, color=WHITE)
add_text(slide, '0,0          se o cliente não tem histórico de pedidos\n'
                '[0,5 ; 1,0]  se tem pedidos: 0,5 + 0,5 × min(1, pedidos/mês ÷ 10)',
         Inches(0.55), Inches(1.95), Inches(5.7), Inches(0.9),
         size=10.5, color=rgb('#E1BEE7'), italic=True)
add_text(slide, '→ Mede QUANTO de informação está disponível para a decisão',
         Inches(0.55), Inches(2.85), Inches(5.7), Inches(0.7),
         size=10.5, color=rgb('#CE93D8'))

# Score B box
add_rect(slide, Inches(0.4), Inches(4.1), Inches(6.0), Inches(2.55), rgb('#1A237E'))
add_text(slide, 'Score B — Complexidade de Perfil',
         Inches(0.55), Inches(4.2), Inches(5.7), Inches(0.45),
         size=13, bold=True, color=WHITE)
add_text(slide, 'score_B = 1 − |p_aplicacao − 0,5| × 2\n'
                '→ 1,0 quando p_app ≈ 0,5 (modelo incerto)\n'
                '→ 0,0 quando p_app ≈ 0 ou 1 (perfil claro)',
         Inches(0.55), Inches(4.65), Inches(5.7), Inches(1.1),
         size=10.5, color=rgb('#BBDEFB'), italic=True)
add_text(slide, '→ Mede QUÃO DIFÍCIL é classificar este cliente',
         Inches(0.55), Inches(5.75), Inches(5.7), Inches(0.7),
         size=10.5, color=rgb('#90CAF9'))

# Rules box
add_rect(slide, Inches(6.7), Inches(1.4), Inches(6.3), Inches(5.25), DARK)
add_text(slide, 'Regras de Roteamento',
         Inches(6.85), Inches(1.5), Inches(6.0), Inches(0.45),
         size=13, bold=True, color=WHITE)

rules = [
    ('BEHAVIORAL',    'Score A ≥ 0,5',                    LBLUE,  '#E3F2FD'),
    ('APPLICATION',   'Score A < 0,5  E  Score B < 0,6',  LORANG, '#FFF3E0'),
    ('MANUAL_REVIEW', 'Score A < 0,5  E  Score B ≥ 0,6',  LPURP,  '#F3E5F5'),
]
for i, (tier, rule, bg, txt_bg) in enumerate(rules):
    y = Inches(2.1 + i * 1.5)
    add_rect(slide, Inches(6.85), y, Inches(5.9), Inches(1.3), bg)
    add_text(slide, tier,
             Inches(7.0), y + Inches(0.1), Inches(5.6), Inches(0.42),
             size=12, bold=True, color=WHITE)
    add_text(slide, 'Condição: ' + rule,
             Inches(7.0), y + Inches(0.55), Inches(5.6), Inches(0.35),
             size=10.5, color=rgb(txt_bg))

add_rect(slide, Inches(6.85), Inches(5.9), Inches(5.9), Inches(0.65), GREEN)
add_text(slide, 'AUC do sistema: 0,850  (+17,2% vs baseline 0,726)',
         Inches(7.0), Inches(5.92), Inches(5.6), Inches(0.6),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Resultados do Router
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide_header('Resultados do Model Router', PURPLE,
                         'Distribuição de roteamento e métricas por tier')
add_slide_number(slide, 10)

add_image(slide, img_routing, Inches(0.4), Inches(1.3), Inches(7.0), Inches(4.2))

# Tier detail boxes
tier_details = [
    ('BEHAVIORAL',    '664 (22,1%)', '0,770', '20,8%', LBLUE),
    ('APPLICATION',   '998 (33,3%)', '0,772', '33,4%', LORANG),
    ('MANUAL_REVIEW', '1.338 (44,6%)', 'App (flagged)', '35,1%', LPURP),
]
for i, (tier, clients, auc, dr, color) in enumerate(tier_details):
    y = Inches(1.35 + i * 1.85)
    add_rect(slide, Inches(7.7), y, Inches(5.4), Inches(1.65), color)
    add_text(slide, tier,
             Inches(7.85), y + Inches(0.08), Inches(5.1), Inches(0.4),
             size=12, bold=True, color=WHITE)
    for j, (label, val) in enumerate([
        ('Clientes', clients), ('Modelo', auc), ('Taxa inadimp.', dr)
    ]):
        x = Inches(7.85 + j * 1.78)
        add_text(slide, label, x, y + Inches(0.52), Inches(1.6), Inches(0.3),
                 size=8.5, color=rgb('#E3F2FD'))
        add_text(slide, val, x, y + Inches(0.82), Inches(1.6), Inches(0.6),
                 size=13, bold=True, color=WHITE)

add_rect(slide, Inches(0.4), Inches(5.7), Inches(12.7), Inches(1.55), DARK)
add_text(slide, 'Impacto no AUC do sistema',
         Inches(0.55), Inches(5.78), Inches(12.4), Inches(0.4),
         size=12, bold=True, color=WHITE)
auc_compare = [
    ('Baseline único modelo', '0,726', MGRAY),
    ('Sistema Rule-based', '0,850', GREEN),
    ('Sistema ML Router', '0,848', TEAL),
    ('Lift (rule-based)', '+4,6pp', ORANGE),
]
for i, (label, val, color) in enumerate(auc_compare):
    x = Inches(0.55 + i * 3.15)
    add_rect(slide, x, Inches(6.22), Inches(2.9), Inches(0.85), color)
    add_text(slide, f'{label}: {val}', x + Inches(0.1), Inches(6.28),
             Inches(2.7), Inches(0.72), size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Interpretabilidade
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide_header('Interpretabilidade — SHAP', MGRAY,
                         'Explicabilidade global e local via SHAP TreeExplainer')
add_slide_number(slide, 11)

add_image(slide, img_shap, Inches(0.4), Inches(1.3), Inches(8.5), Inches(4.5))

add_rect(slide, Inches(9.1), Inches(1.35), Inches(4.0), Inches(5.9), DARK)
add_text(slide, 'O que o SHAP revela?',
         Inches(9.25), Inches(1.45), Inches(3.7), Inches(0.42),
         size=12, bold=True, color=WHITE)
add_bullets(slide, [
    'CNPJ novo → maior risco\n(idade_cnpj_mid é o driver #1 em ambos)',
    '',
    'Sócio com negativação no\nSerasa amplifica muito o risco',
    '',
    'Diversificação de credores\n(n_setores) indica fragilidade',
    '',
    'Comportamento de pagamento\n(delay_mean) é top-2 no\nModelo Comportamental',
    '',
    'Curva de calibração validada:\npredições alinhadas com\ntaxas reais observadas',
], Inches(9.25), Inches(1.9), Inches(3.7), Inches(5.1),
   size=10, color=rgb('#B0BEC5'), bullet='▸ ', spacing=5)

add_rect(slide, Inches(0.4), Inches(6.0), Inches(8.5), Inches(1.25), rgb('#263238'))
add_text(slide, 'Verificação por permutação: ranking de importância consistente com SHAP em ambos os modelos.',
         Inches(0.55), Inches(6.08), Inches(8.2), Inches(0.45),
         size=10.5, bold=True, color=rgb('#FFF176'))
add_text(slide, 'Model cards exportados com metadados, thresholds, métricas e fontes de dados de cada modelo.',
         Inches(0.55), Inches(6.52), Inches(8.2), Inches(0.65),
         size=10, color=rgb('#B0BEC5'))

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Política de Crédito
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide_header('Política de Crédito', GREEN,
                         'Thresholds por tier · Limites de crédito baseados em risco')
add_slide_number(slide, 12)

# Decision rules table
headers = ['Tier', 'Auto-Aprovar (p <)', 'Manual (faixa)', 'Auto-Negar (p ≥)']
h_widths = [2.5, 2.5, 2.6, 2.5]
h_colors = [DARK, GREEN, AMBER, RED]
start_x = Inches(0.4)
row_y   = Inches(1.4)

for j, (h, w, c) in enumerate(zip(headers, h_widths, h_colors)):
    x = start_x + sum(Inches(w2) for w2 in h_widths[:j])
    add_rect(slide, x, row_y, Inches(w), Inches(0.5), c)
    add_text(slide, h, x + Inches(0.08), row_y + Inches(0.05),
             Inches(w - 0.15), Inches(0.4), size=11, bold=True, color=WHITE)

rows = [
    ('BEHAVIORAL',    '0,25', '0,25 – 0,35', '0,35', LBLUE),
    ('APPLICATION',   '0,15', '0,15 – 0,25', '0,25', LORANG),
    ('MANUAL_REVIEW', '0,40*', '—', '0,40*', LPURP),
]
for i, (tier, a, m, d, color) in enumerate(rows):
    ry = row_y + Inches(0.5 + i * 0.6)
    bg = rgb('#F5F5F5') if i % 2 == 0 else WHITE
    vals = [tier, a, m, d]
    for j, (val, w, hc) in enumerate(zip(vals, h_widths, h_colors)):
        x = start_x + sum(Inches(w2) for w2 in h_widths[:j])
        c2 = color if j == 0 else bg
        add_rect(slide, x, ry, Inches(w), Inches(0.55), c2)
        add_text(slide, val, x + Inches(0.08), ry + Inches(0.08),
                 Inches(w - 0.15), Inches(0.4), size=11,
                 bold=(j == 0), color=WHITE if j == 0 else DARK)

add_text(slide, '* Threshold indicativo para analista humano no MANUAL_REVIEW',
         Inches(0.4), Inches(3.35), Inches(10), Inches(0.3),
         size=8.5, italic=True, color=MGRAY)

# Credit limits
add_rect(slide, Inches(0.4), Inches(3.7), Inches(12.7), Inches(0.4), DARK)
add_text(slide, 'Limites de Crédito por Tier e Faixa de Risco (p)',
         Inches(0.55), Inches(3.72), Inches(12.4), Inches(0.36),
         size=12, bold=True, color=WHITE)

limit_data = [
    ('BEHAVIORAL', [('p < 0,10', 'R$25.000'), ('0,10 – 0,18', 'R$15.000'), ('0,18 – 0,25', 'R$10.000')], LBLUE),
    ('APPLICATION', [('p < 0,08', 'R$20.000'), ('0,08 – 0,12', 'R$12.000'), ('0,12 – 0,15', 'R$8.000')], LORANG),
    ('MANUAL_REVIEW', [('p < 0,40', 'R$8.000*')], LPURP),
]

for i, (tier, bands, color) in enumerate(limit_data):
    x = Inches(0.4 + i * 4.27)
    add_rect(slide, x, Inches(4.15), Inches(4.0), Inches(0.4), color)
    add_text(slide, tier, x + Inches(0.12), Inches(4.18), Inches(3.75), Inches(0.35),
             size=11, bold=True, color=WHITE)
    for k, (band, limit) in enumerate(bands):
        y = Inches(4.6 + k * 0.7)
        add_rect(slide, x, y, Inches(4.0), Inches(0.62), rgb('#ECEFF1'))
        add_text(slide, band, x + Inches(0.12), y + Inches(0.08),
                 Inches(2.2), Inches(0.45), size=10.5, color=DARK)
        add_text(slide, limit, x + Inches(2.3), y + Inches(0.08),
                 Inches(1.6), Inches(0.45), size=11, bold=True, color=color,
                 align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — Análise de Cenários
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide_header('Análise de Cenários', GREEN,
                         'Comparação de quatro políticas de crédito no mesmo portfólio de 3.000 clientes')
add_slide_number(slide, 13)

add_image(slide, img_scenarios, Inches(0.4), Inches(1.35), Inches(7.9), Inches(4.0))

# Scenario detail
add_rect(slide, Inches(8.5), Inches(1.35), Inches(4.6), Inches(5.9), DARK)
add_text(slide, 'Detalhamento',
         Inches(8.65), Inches(1.45), Inches(4.3), Inches(0.4),
         size=12, bold=True, color=WHITE)

scenario_details = [
    ('Baseline', '21,4%', '5,6%', 'R$0,43M', MGRAY),
    ('Recomendado', '39,0%', '7,3%', 'R$1,03M', PURPLE),
    ('Conservador', '21,4%', '2,8%', 'R$0,22M', BLUE),
    ('Permissivo', '60,7%', '11,7%', 'R$2,56M', RED),
]
hdrs = ['Cenário', 'Aprovação', 'Inadimp.', 'Perda esp.']
hs   = [1.45, 1.0, 1.0, 1.0]
xpos = [Inches(8.65 + sum(Inches(h) for h in hs[:j])) for j in range(4)]
ry   = Inches(1.9)
for j, (hdr, hc) in enumerate(zip(hdrs, [DARK]*4)):
    add_rect(slide, xpos[j], ry, Inches(hs[j]), Inches(0.45), MGRAY)
    add_text(slide, hdr, xpos[j] + Inches(0.05), ry + Inches(0.05),
             Inches(hs[j] - 0.08), Inches(0.35), size=9, bold=True, color=WHITE)

for i, (name, app, dr, loss, color) in enumerate(scenario_details):
    ry2 = Inches(2.4 + i * 0.95)
    bg  = color if name == 'Recomendado' else (rgb('#F5F5F5') if i % 2 == 0 else WHITE)
    for j, (val, h) in enumerate(zip([name, app, dr, loss], hs)):
        add_rect(slide, xpos[j], ry2, Inches(h), Inches(0.85),
                 color if j == 0 else bg)
        add_text(slide, val, xpos[j] + Inches(0.05), ry2 + Inches(0.18),
                 Inches(h - 0.08), Inches(0.55), size=10.5,
                 bold=(j == 0 or name == 'Recomendado'),
                 color=WHITE if (j == 0 or name == 'Recomendado') else DARK,
                 align=PP_ALIGN.CENTER)

add_rect(slide, Inches(8.5), Inches(5.55), Inches(4.6), Inches(1.7), GREEN)
add_text(slide, '★  Política Recomendada',
         Inches(8.65), Inches(5.62), Inches(4.3), Inches(0.42),
         size=12, bold=True, color=WHITE)
add_bullets(slide, [
    'Quase 2× mais aprovações que o baseline',
    'Aumento de inadimplência aceitável (+1,7pp)',
    'Clientes BEHAVIORAL absorvem limites maiores',
    'com menor risco real (20,8% vs 33,4%)',
], Inches(8.65), Inches(6.05), Inches(4.3), Inches(1.2),
   size=9.5, color=rgb('#C8E6C9'), bullet='▸ ', spacing=4)

add_rect(slide, Inches(0.4), Inches(5.55), Inches(7.9), Inches(1.7), rgb('#263238'))
add_text(slide, 'Premissa: limite uniforme de R$12.000 para comparação entre cenários. '
         'Limites reais por faixa de risco aumentam a diferença de performance entre o router e o baseline.',
         Inches(0.55), Inches(5.62), Inches(7.6), Inches(1.6),
         size=10.5, italic=True, color=rgb('#90A4AE'))

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — KPIs Consolidados
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide_header('Métricas Consolidadas', DARK,
                         'Resultados finais do sistema de crédito inteligente com Model Routing')
add_slide_number(slide, 14)

kpis = [
    ('0,850',  'AUC do Sistema\n(Rule-based Router)', PURPLE),
    ('+4,6pp', 'Lift vs.\nBaseline único', GREEN),
    ('39,0%',  'Taxa de\nAprovação', TEAL),
    ('3 tiers','Roteamento\nAutomático', BLUE),
    ('R$9,6M', 'Exposição\nTotal (crédito)', ORANGE),
    ('16,7%',  'Taxa de Perda\nEsperada', RED),
]

for i, (val, lbl, color) in enumerate(kpis):
    row = i // 3
    col = i  % 3
    x = Inches(0.5 + col * 4.27)
    y = Inches(1.5 + row * 2.6)
    add_rect(slide, x, y, Inches(4.0), Inches(2.3), color)
    add_text(slide, val,
             x, y + Inches(0.25), Inches(4.0), Inches(1.1),
             size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, lbl,
             x, y + Inches(1.35), Inches(4.0), Inches(0.8),
             size=12, color=rgb('#E3F2FD'), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — Governança & Próximos Passos
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide_header('Governança & Próximos Passos', DARK,
                         'Monitoramento contínuo, atualização de modelos e evolução da arquitetura')
add_slide_number(slide, 15)

gov_sections = [
    ('Monitoramento', [
        'Mensal: taxa de inadimplência por tier vs. esperada (alerta se Δ > 5%)',
        'Trimestral: AUC em coortes com 90 dias de maturidade',
        'Trimestral: representação por segmento no tier MANUAL_REVIEW (viés)',
        'Anual: retreino completo com dados atualizados',
    ], DARK),
    ('Gatilhos de Atualização', [
        'AUC cai abaixo de 0,70 em coorte de 90 dias',
        'Taxa de inadimplência dos AUTO-APPROVE ultrapassa 15% (BEHAVIORAL) ou 10% (APPLICATION)',
        'Mudança significativa no mix de produtos ou segmento de clientes',
        'Evento macroeconômico disruptivo (ex: variação >20% na inadimplência do setor)',
    ], RED),
    ('Próximos Passos Técnicos', [
        'Integrar src/router.py na API de originação de crédito',
        'Construir dashboard de monitoramento em tempo real',
        'Expandir tier BEHAVIORAL: coletar histórico dos 77,9% sem dados (meta: 40%+)',
        'Avaliar modelos de séries temporais para capturar sazonalidade nos pedidos',
        'Testar LightGBM e CatBoost no Modelo de Aplicação',
    ], TEAL),
]

for i, (title, items, color) in enumerate(gov_sections):
    x = Inches(0.4 + i * 4.27)
    add_rect(slide, x, Inches(1.4), Inches(4.0), Inches(5.9), color)
    add_text(slide, title, x + Inches(0.15), Inches(1.5), Inches(3.7), Inches(0.45),
             size=13, bold=True, color=WHITE)
    add_bullets(slide, items, x + Inches(0.15), Inches(2.0), Inches(3.7), Inches(5.2),
                size=10, color=WHITE, bullet='▸ ', spacing=8)

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUT_PATH)
print(f'Apresentação salva → {OUT_PATH}')
print(f'Total de slides: {len(prs.slides)}')
